"""Build and render fixed, traceable scenario reports.

The assembler deliberately delegates DC-plan metrics to ``AnalyticsService``;
the report is a presentation of existing analytics, not a second metric SQL
implementation.
"""

from __future__ import annotations

import base64
import io
import os
import subprocess
import sys
from datetime import datetime, timezone
from pathlib import Path

from planalign_api.models.report import (
    ReportProvenance,
    ReportWarning,
    ReportYearMetrics,
    ScenarioReport,
)
from planalign_api.services.analytics_service import AnalyticsService
from planalign_api.services.database_path_resolver import (
    create_api_database_path_resolver,
)
from planalign_api.services.run_trust import add_current_config_drift, read_run_trust
from planalign_api.storage.workspace_storage import WorkspaceStorage


class ReportNotFoundError(ValueError):
    """The requested scenario has no usable result."""


def build_scenario_report(
    storage: WorkspaceStorage,
    workspace_id: str,
    scenario_id: str,
    comparison_ids: list[str] | None = None,
) -> ScenarioReport:
    """Assemble one or more scenario series from selected result databases."""
    scenario = _resolve_scenario(storage, workspace_id, scenario_id)
    if scenario is None:
        raise ReportNotFoundError(f"Scenario {scenario_id} was not found")
    selected = [scenario_id, *(comparison_ids or [])]
    report = ScenarioReport(
        title=f"PlanAlign scenario report: {scenario.name}",
        generated_at=datetime.now(timezone.utc),
        comparison_scenario_ids=selected[1:],
    )
    for index, selected_id in enumerate(selected):
        selected_scenario = _resolve_scenario(storage, workspace_id, selected_id)
        if selected_scenario is None:
            raise ReportNotFoundError(f"Scenario {selected_id} was not found")
        if selected_scenario.status != "completed":
            report.warnings.append(
                ReportWarning(
                    code="run_not_completed",
                    message=f"Scenario {selected_id} is {selected_scenario.status}; no authoritative report was produced.",
                    severity="error",
                )
            )
            continue
        resolved = create_api_database_path_resolver(storage).resolve(
            workspace_id, selected_id, verify_database=False
        )
        if not resolved.exists or resolved.path is None:
            report.warnings.append(
                ReportWarning(
                    code="result_unavailable",
                    message=f"Selected result for scenario {selected_id} is unavailable.",
                    severity="error",
                )
            )
            continue
        metrics = _append_scenario(
            report, storage, workspace_id, selected_scenario, resolved.path
        )
        if index > 0 and metrics:
            report.comparison_years[selected_scenario.id] = metrics
    if not report.provenance and not report.warnings:
        raise ReportNotFoundError("No completed selected results were found")
    return report


def _append_scenario(
    report, storage, workspace_id, scenario, database_path
) -> list[ReportYearMetrics]:
    analytics = AnalyticsService(storage).get_dc_plan_analytics(
        workspace_id, scenario.id, scenario.name
    )
    if analytics is None:
        report.warnings.append(
            ReportWarning(
                code="analytics_unavailable",
                message=f"Analytics are unavailable for scenario {scenario.id}.",
                severity="error",
            )
        )
        return []
    workforce = _workforce_by_year(database_path)
    contribution = {item.year: item for item in analytics.contribution_by_year}
    years = sorted(set(workforce) | set(contribution))
    if not years:
        report.warnings.append(
            ReportWarning(
                code="metrics_unavailable",
                message=f"No yearly metrics are available for scenario {scenario.id}.",
                severity="error",
            )
        )
        return []
    metrics = [
        ReportYearMetrics(
            year=year,
            headcount=workforce.get(year, {}).get("headcount"),
            average_compensation=workforce.get(year, {}).get("average_compensation"),
            participation_rate=getattr(
                contribution.get(year), "participation_rate", None
            ),
            employer_cost=(
                getattr(contribution.get(year), "total_employer_match", 0.0)
                + getattr(contribution.get(year), "total_employer_core", 0.0)
            )
            if contribution.get(year)
            else None,
        )
        for year in years
    ]
    if not report.years:
        report.years = metrics
    trust = add_current_config_drift(
        read_run_trust(database_path, scenario.last_run_id),
        storage.get_merged_config(workspace_id, scenario.id),
    )
    warnings = _trust_warnings(trust)
    report.warnings.extend(warnings)
    report.provenance.append(
        ReportProvenance(
            workspace_id=workspace_id,
            scenario_id=scenario.id,
            scenario_name=scenario.name,
            run_id=trust.run_id or scenario.last_run_id,
            config_fingerprint=trust.config_fingerprint,
            random_seed=trust.random_seed,
            run_timestamp=trust.run_timestamp,
            source_result=str(database_path),
            git_sha=_git_sha(),
        )
    )
    report.assumptions.setdefault(
        "config", storage.get_merged_config(workspace_id, scenario.id) or {}
    )
    return metrics


def _resolve_scenario(storage: WorkspaceStorage, workspace_id: str, identifier: str):
    """Accept stable IDs and human-facing scenario names in CLI workflows."""
    scenario = storage.get_scenario(workspace_id, identifier)
    if scenario is not None:
        return scenario
    return next(
        (
            item
            for item in storage.list_scenarios(workspace_id)
            if item.name == identifier
        ),
        None,
    )


def _workforce_by_year(database_path: Path) -> dict[int, dict[str, float | int]]:
    import duckdb

    query = """
        SELECT simulation_year,
               COUNT(*) AS headcount,
               ROUND(AVG(current_compensation), 2) AS average_compensation
        FROM fct_workforce_snapshot
        GROUP BY simulation_year ORDER BY simulation_year
    """
    with duckdb.connect(str(database_path), read_only=True) as conn:
        return {
            int(year): {
                "headcount": int(headcount),
                "average_compensation": float(avg or 0),
            }
            for year, headcount, avg in conn.execute(query).fetchall()
        }


def _trust_warnings(trust) -> list[ReportWarning]:
    return [
        ReportWarning(
            code=reason,
            message={
                "mixed_generation": "The selected result contains mixed configuration or seed generations.",
                "current_config_mismatch": "Current configuration differs from the selected result.",
                "current_seed_mismatch": "Current seed differs from the selected result.",
            }.get(reason, reason),
        )
        for reason in trust.reasons
    ]


def _git_sha() -> str | None:
    try:
        return subprocess.run(
            ["git", "rev-parse", "HEAD"], capture_output=True, text=True, check=True
        ).stdout.strip()
    except (OSError, subprocess.CalledProcessError):
        return None


def render_report(
    report: ScenarioReport, output: Path, fmt: str, template_dir: Path | None = None
) -> Path:
    """Render HTML, PDF, or PPTX; optional render dependencies fail explicitly."""
    output.parent.mkdir(parents=True, exist_ok=True)
    html = _render_html(report, template_dir)
    if fmt == "html":
        output.write_text(html, encoding="utf-8")
    elif fmt == "pdf":
        _configure_macos_weasyprint_paths()
        try:
            from weasyprint import HTML
        except (ImportError, OSError) as exc:
            raise RuntimeError(
                "PDF reports require WeasyPrint plus its native macOS libraries. "
                "Install with 'brew install weasyprint', then restart Studio."
            ) from exc
        try:
            HTML(string=html, base_url=str(output.parent)).write_pdf(str(output))
        except (OSError, AssertionError) as exc:
            raise RuntimeError(
                "WeasyPrint could not render the report template. "
                "Restart Studio after updating PlanAlign and retry."
            ) from exc
    elif fmt == "pptx":
        _render_pptx(report, output)
    else:
        raise ValueError(f"Unsupported report format: {fmt}")
    return output


def _configure_macos_weasyprint_paths() -> None:
    """Expose Homebrew's native libraries to WeasyPrint on macOS."""
    if sys.platform != "darwin":
        return
    candidates = [
        "/opt/homebrew/lib",
        "/opt/homebrew/opt/glib/lib",
        "/opt/homebrew/opt/pango/lib",
        "/opt/homebrew/opt/cairo/lib",
        "/usr/local/lib",
        "/usr/local/opt/glib/lib",
        "/usr/local/opt/pango/lib",
        "/usr/local/opt/cairo/lib",
    ]
    existing = [path for path in candidates if Path(path).is_dir()]
    if existing:
        current = os.environ.get("DYLD_FALLBACK_LIBRARY_PATH", "")
        os.environ["DYLD_FALLBACK_LIBRARY_PATH"] = ":".join(
            dict.fromkeys(existing + ([current] if current else []))
        )


def _render_html(report: ScenarioReport, template_dir: Path | None) -> str:
    rows = "".join(
        f"<tr><td>{m.year}</td><td>{m.headcount if m.headcount is not None else 'Unavailable'}</td>"
        f"<td>{_money(m.average_compensation)}</td><td>{_pct(m.participation_rate)}</td>"
        f"<td>{_money(m.employer_cost)}</td></tr>"
        for m in report.years
    )
    provenance = "".join(
        f"<li>{p.scenario_name} ({p.scenario_id}) — run {p.run_id or 'unavailable'}, "
        f"config {p.config_fingerprint or 'unavailable'}, seed {p.random_seed if p.random_seed is not None else 'unavailable'}, "
        f"git {p.git_sha or 'unavailable'}</li>"
        for p in report.provenance
    )
    warnings = "".join(
        f"<li class='{w.severity}'>{w.message}</li>" for w in report.warnings
    )
    charts = "".join(
        f"<img src='{_chart_data(report.years, field, label)}' alt='{label}'/>"
        for field, label in (
            ("headcount", "Headcount"),
            ("average_compensation", "Average compensation"),
            ("participation_rate", "Participation"),
            ("employer_cost", "Employer cost"),
        )
    )
    trace = " | ".join(
        f"{p.scenario_id} / config {p.config_fingerprint or 'unavailable'}"
        for p in report.provenance
    )
    template = (
        (template_dir / "report.html").read_text(encoding="utf-8")
        if template_dir
        else _DEFAULT_TEMPLATE
    )
    return (
        template.replace("{{ title }}", report.title)
        .replace("{{ rows }}", rows)
        .replace("{{ provenance }}", provenance)
        .replace("{{ warnings }}", warnings)
        .replace("{{ charts }}", charts)
        .replace("{{ trace }}", trace)
    )


def _chart_data(metrics: list[ReportYearMetrics], field: str, label: str) -> str:
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    values = [getattr(item, field) for item in metrics]
    years = [item.year for item in metrics]
    figure, axis = plt.subplots(figsize=(4.2, 2.1), dpi=120)
    axis.plot(
        years,
        [value if value is not None else float("nan") for value in values],
        color="#2f7d32",
        marker="o",
    )
    axis.set_title(label, fontsize=9)
    axis.grid(alpha=0.2)
    figure.tight_layout()
    stream = io.BytesIO()
    figure.savefig(stream, format="png", transparent=False)
    plt.close(figure)
    return "data:image/png;base64," + base64.b64encode(stream.getvalue()).decode(
        "ascii"
    )


def _render_pptx(report: ScenarioReport, output: Path) -> None:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "PowerPoint reports require the optional python-pptx dependency"
        ) from exc
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = report.title
    box = slide.shapes.add_textbox(1e6, 1.5e6, 8e6, 5e6)
    text = "Year | Headcount | Avg comp | Participation | Employer cost\n"
    text += "\n".join(
        f"{m.year} | {m.headcount or 'N/A'} | {_money(m.average_compensation)} | {_pct(m.participation_rate)} | {_money(m.employer_cost)}"
        for m in report.years
    )
    box.text_frame.text = text
    _add_pptx_footer(slide, report)
    chart_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    chart_slide.shapes.title.text = "Standard metrics"
    for index, (field, label) in enumerate(
        (
            ("headcount", "Headcount"),
            ("average_compensation", "Average compensation"),
            ("participation_rate", "Participation"),
            ("employer_cost", "Employer cost"),
        )
    ):
        raw = base64.b64decode(_chart_data(report.years, field, label).split(",", 1)[1])
        chart_slide.shapes.add_picture(
            io.BytesIO(raw),
            400000 + (index % 2) * 4500000,
            1100000 + (index // 2) * 2600000,
            width=4000000,
        )
    _add_pptx_footer(chart_slide, report)
    provenance = presentation.slides.add_slide(presentation.slide_layouts[5])
    provenance.shapes.title.text = "Assumptions and provenance"
    provenance.shapes.add_textbox(1e6, 1.5e6, 8e6, 5e6).text_frame.text = "\n".join(
        f"{item.scenario_id}: run={item.run_id or 'unavailable'} config={item.config_fingerprint or 'unavailable'} git={item.git_sha or 'unavailable'}"
        for item in report.provenance
    )
    _add_pptx_footer(provenance, report)
    presentation.save(output)


def _add_pptx_footer(slide, report: ScenarioReport) -> None:
    trace = " | ".join(
        f"{item.scenario_id} / config {item.config_fingerprint or 'unavailable'}"
        for item in report.provenance
    )
    slide.shapes.add_textbox(
        400000, 6800000, 9000000, 250000
    ).text_frame.text = f"Traceable PlanAlign report — {trace}"


def _money(value: float | None) -> str:
    return "Unavailable" if value is None else f"${value:,.0f}"


def _pct(value: float | None) -> str:
    return "Unavailable" if value is None else f"{value:.1f}%"


_DEFAULT_TEMPLATE = """<!doctype html><html><head><meta charset='utf-8'><style>@page{size:letter landscape;margin:0.45in}body{font-family:Arial;color:#273238}h1{color:#2f7d32;border-bottom:4px solid #2f7d32;padding-bottom:12px}table{border-collapse:collapse;width:100%;margin-top:18px}th{background:#2f7d32;color:white}td,th{padding:9px;border:1px solid #d6ded9;text-align:right}td:first-child,th:first-child{text-align:left}.charts{display:block;width:100%}.charts img{display:inline-block;width:46%;vertical-align:top;margin:6px}.warning{color:#9a6700}.error{color:#b42318}footer{position:fixed;bottom:0;font-size:9px;color:#64736d}</style></head><body><h1>{{ title }}</h1><p>Executive summary generated from the selected PlanAlign result.</p><table><thead><tr><th>Year</th><th>Headcount</th><th>Average compensation</th><th>Participation</th><th>Employer cost</th></tr></thead><tbody>{{ rows }}</tbody></table><div class='charts'>{{ charts }}</div><h2>Assumptions and provenance</h2><ul>{{ provenance }}</ul><ul>{{ warnings }}</ul><footer>Traceable PlanAlign report — {{ trace }}</footer></body></html>"""
